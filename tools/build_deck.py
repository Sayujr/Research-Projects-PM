"""Build a 6-slide presentation deck (.pptx) for the research associate.

Run:   .venv/bin/python tools/build_deck.py [output_filename.pptx]
  (default filename: PRESENTATION.pptx at repo root)

Design: corporate-clean, 16:9, Calibri, shapes (not ASCII/Mermaid), editable.
"""

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


# ---------- palette ----------

NAVY    = RGBColor(0x1a, 0x23, 0x32)
TEAL    = RGBColor(0x2d, 0x7a, 0x7a)
OFFWH   = RGBColor(0xfa, 0xfa, 0xf9)
INK     = RGBColor(0x2b, 0x33, 0x40)
MUTED   = RGBColor(0x6b, 0x72, 0x80)
LINE    = RGBColor(0xd1, 0xd5, 0xdb)
CARD_BG = RGBColor(0xff, 0xff, 0xff)

DEEP    = RGBColor(0x25, 0x63, 0xeb)   # writing / analysis
LAB     = RGBColor(0x8b, 0x5c, 0xf6)   # bench work
LIGHT   = RGBColor(0x6b, 0x72, 0x80)   # admin / light
MEETING = RGBColor(0xea, 0x58, 0x0c)   # social / meeting
PERS    = RGBColor(0x64, 0x74, 0x8b)   # personal

GREEN   = RGBColor(0x2d, 0xa7, 0x4a)
AMBER   = RGBColor(0xf0, 0xb4, 0x29)
RED     = RGBColor(0xdc, 0x26, 0x26)

WHITE   = RGBColor(0xff, 0xff, 0xff)

# ---------- canvas ----------

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.5)

FONT = "Calibri"

DECK_TITLE = "research-pm"


# ---------- helpers ----------

def tint(rgb: RGBColor, pct: float) -> RGBColor:
    """Lighten an RGB toward white by pct (0..1)."""
    r, g, b = rgb[0], rgb[1], rgb[2]
    return RGBColor(
        int(r + (255 - r) * pct),
        int(g + (255 - g) * pct),
        int(b + (255 - b) * pct),
    )


def set_solid_fill(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def set_line(shape, color=None, width=None):
    if color is None:
        shape.line.fill.background()
        return
    shape.line.color.rgb = color
    if width is not None:
        shape.line.width = width


def set_line_dashed(shape, color: RGBColor, width_pt: float = 1.25):
    """Apply a dashed stroke via raw XML (python-pptx has no first-class API)."""
    shape.line.color.rgb = color
    shape.line.width = Pt(width_pt)
    ln = shape.line._get_or_add_ln()
    for tag in ("a:prstDash",):
        for el in ln.findall(qn(tag)):
            ln.remove(el)
    prstDash = etree.SubElement(ln, qn("a:prstDash"))
    prstDash.set("val", "dash")


def add_textbox(slide, x, y, w, h, text, *,
                size=14, bold=False, color=INK, font=FONT,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_rect(slide, x, y, w, h, *, fill=None, line_color=None,
             line_width=None, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        s.fill.background()
    else:
        set_solid_fill(s, fill)
    set_line(s, line_color, line_width)
    s.shadow.inherit = False
    # blank any auto text
    if s.has_text_frame:
        s.text_frame.text = ""
    return s


def add_rrect(slide, x, y, w, h, *, fill=None, line_color=None, line_width=None):
    """Rounded rectangle with a gentle corner radius."""
    s = add_rect(slide, x, y, w, h, fill=fill, line_color=line_color,
                 line_width=line_width, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    # reduce corner radius (default is ~17% of short side; make it subtler)
    try:
        s.adjustments[0] = 0.08
    except Exception:
        pass
    return s


def add_line(slide, x1, y1, x2, y2, *, color=LINE, width_pt=1.0):
    s = slide.shapes.add_connector(1, x1, y1, x2, y2)  # 1 = straight
    s.line.color.rgb = color
    s.line.width = Pt(width_pt)
    return s


def add_dashed_line(slide, x1, y1, x2, y2, *, color=TEAL, width_pt=1.25):
    s = slide.shapes.add_connector(1, x1, y1, x2, y2)
    set_line_dashed(s, color, width_pt)
    return s


def add_arrow(slide, x1, y1, x2, y2, *, color=MUTED, width_pt=1.25):
    s = slide.shapes.add_connector(1, x1, y1, x2, y2)
    s.line.color.rgb = color
    s.line.width = Pt(width_pt)
    # add arrow head at tail
    ln = s.line._get_or_add_ln()
    for tag in ("a:headEnd", "a:tailEnd"):
        for el in ln.findall(qn(tag)):
            ln.remove(el)
    tailEnd = etree.SubElement(ln, qn("a:tailEnd"))
    tailEnd.set("type", "triangle")
    tailEnd.set("w", "med")
    tailEnd.set("h", "med")
    return s


def add_dot(slide, cx, cy, d, color):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - d/2, cy - d/2, d, d)
    set_solid_fill(s, color)
    set_line(s, None)
    return s


# ---------- chrome (per-slide title + footer) ----------

def add_chrome(slide, slide_num, title_text, eyebrow=None):
    # off-white background
    bg = add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=OFFWH)
    # leave it in the back
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)

    # eyebrow (small caps label above title)
    if eyebrow:
        add_textbox(slide, MARGIN, Inches(0.35),
                    Inches(12.33), Inches(0.25),
                    eyebrow.upper(), size=9, bold=True, color=TEAL,
                    line_spacing=1.0)
        title_y = Inches(0.6)
    else:
        title_y = Inches(0.45)

    # title
    add_textbox(slide, MARGIN, title_y,
                Inches(11), Inches(0.7),
                title_text, size=28, bold=True, color=NAVY,
                line_spacing=1.0)

    # teal accent stripe under title
    add_rect(slide, MARGIN, Inches(1.22),
             Inches(0.8), Inches(0.05),
             fill=TEAL, line_color=None)

    # footer line
    add_line(slide, MARGIN, Inches(7.05),
             SLIDE_W - MARGIN, Inches(7.05),
             color=LINE, width_pt=0.75)

    # footer text — left deck title, right slide num
    add_textbox(slide, MARGIN, Inches(7.12),
                Inches(6), Inches(0.25),
                DECK_TITLE, size=9, color=MUTED, line_spacing=1.0)
    add_textbox(slide, SLIDE_W - MARGIN - Inches(2), Inches(7.12),
                Inches(2), Inches(0.25),
                f"{slide_num} / 6", size=9, color=MUTED,
                align=PP_ALIGN.RIGHT, line_spacing=1.0)


# ---------- slide 1: Her week with Claude in it ----------

def build_slide_1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_chrome(slide, 1, "A week with this tool",
               eyebrow="Workflow")

    # intro line
    add_textbox(slide, MARGIN, Inches(1.42),
                Inches(12.33), Inches(0.5),
                "This is a markdown-based project manager that Claude operates. You do the work; Claude handles the planning, logging, and follow-ups. The small tags below are the commands Claude runs — a standup (today's plan), a digest (what shipped), a retrospective (what slipped).",
                size=12, color=MUTED, line_spacing=1.2)

    # 5 day columns
    days = ["MON", "TUE", "WED", "THU", "FRI"]
    col_x0 = MARGIN + Inches(0.2)
    col_w = Inches(2.3)
    col_gap = Inches(0.15)
    col_top = Inches(2.0)
    col_bot = Inches(5.25)  # end of task chips band

    # band background (very subtle)
    add_rect(slide, MARGIN, Inches(1.95),
             Inches(12.33), Inches(3.4),
             fill=tint(LINE, 0.6), line_color=None)

    # day headers
    for i, d in enumerate(days):
        x = col_x0 + (col_w + col_gap) * i
        add_textbox(slide, x, col_top,
                    col_w, Inches(0.3),
                    d, size=11, bold=True, color=NAVY,
                    align=PP_ALIGN.CENTER, line_spacing=1.0)

    # task chips: (day_idx, y, h, type, label, duration/dot-color)
    # block types: DEEP, LAB, LIGHT, MEETING
    chips = [
        # Mon
        (0, 2.45, 0.45, MEETING, "/week-plan",        None),
        (0, 2.95, 0.55, DEEP,    "Methods — intro",   "2h · TNF"),
        (0, 3.55, 0.40, LIGHT,   "Admin + email",     "30m"),
        (0, 4.00, 0.45, LAB,     "MCF-7 passage",     "45m"),
        # Tue
        (1, 2.45, 0.55, DEEP,    "Methods draft",     "2h · TNF"),
        (1, 3.05, 0.50, LAB,     "Dose-response",     "2h · TNF"),
        (1, 3.60, 0.45, MEETING, "1:1 Ravi",          "30m"),
        (1, 4.10, 0.40, LIGHT,   "1:1 prep (Sara)",   "25m"),
        # Wed
        (2, 2.45, 0.55, DEEP,    "Results draft",     "2.5h · TNF"),
        (2, 3.05, 0.60, LAB,     "Organoid media",    "3h · Org"),
        (2, 3.70, 0.45, MEETING, "Lab meeting",       "1h"),
        (2, 4.20, 0.40, LIGHT,   "Notes + log",       "20m"),
        # Thu
        (3, 2.45, 0.50, DEEP,    "Methods revise",    "2h · TNF"),
        (3, 3.00, 0.50, LAB,     "Assay prep",        "1.5h · Org"),
        (3, 3.55, 0.45, MEETING, "1:1 Sara",          "30m"),
        (3, 4.05, 0.40, LIGHT,   "Reagent order",     "15m"),
        # Fri
        (4, 2.45, 0.50, DEEP,    "Intro polish",      "2h · TNF"),
        (4, 3.00, 0.45, LAB,     "MCF-7 passage",     "45m"),
        (4, 3.50, 0.45, MEETING, "/digest + /retro",  "30m"),
        (4, 4.00, 0.40, LIGHT,   "Plan next week",    "20m"),
    ]

    for day_idx, y, h, color, label, meta in chips:
        x = col_x0 + (col_w + col_gap) * day_idx + Inches(0.1)
        w = col_w - Inches(0.2)
        chip = add_rrect(slide, x, Inches(y),
                         w, Inches(h),
                         fill=color, line_color=None)
        # label
        add_textbox(slide, x + Inches(0.12), Inches(y + 0.04),
                    w - Inches(0.24), Inches(0.3),
                    label, size=10, bold=True, color=WHITE,
                    line_spacing=1.0)
        if meta:
            add_textbox(slide, x + Inches(0.12), Inches(y + h - 0.28),
                        w - Inches(0.24), Inches(0.25),
                        meta, size=8, color=WHITE, line_spacing=1.0)

    # divider between chip band and Claude-action band
    add_line(slide, MARGIN + Inches(0.2), Inches(5.5),
             SLIDE_W - MARGIN - Inches(0.2), Inches(5.5),
             color=LINE, width_pt=0.75)

    # Claude-action callouts under each day
    add_textbox(slide, MARGIN, Inches(5.6),
                Inches(3), Inches(0.25),
                "CLAUDE RUNS", size=9, bold=True, color=TEAL,
                line_spacing=1.0)

    actions = [
        (0, "Monday",    "/week-plan",  "Drafts your top-5 for the\nweek. You edit."),
        (1, "Morning",   "/standup",    "Today's 3 tasks, anything\noverdue, calendar shape."),
        (2, "Morning",   "/standup",    "Flags anything that slipped\nfrom the day before."),
        (3, "Morning",   "/standup",    "Drafts the 1:1 agenda for\nyour afternoon meeting."),
        (4, "Friday",    "/digest + /retro", "Summary of what shipped,\nnotes on what slipped."),
    ]
    for day_idx, when, cmd, body in actions:
        x = col_x0 + (col_w + col_gap) * day_idx + Inches(0.05)
        w = col_w - Inches(0.1)
        # when label
        add_textbox(slide, x, Inches(5.85),
                    w, Inches(0.22),
                    when.upper(), size=8, bold=True, color=MUTED,
                    line_spacing=1.0)
        # command in teal mono-ish
        add_textbox(slide, x, Inches(6.08),
                    w, Inches(0.28),
                    cmd, size=11, bold=True, color=TEAL,
                    line_spacing=1.0, font="Consolas")
        # body
        add_textbox(slide, x, Inches(6.38),
                    w, Inches(0.6),
                    body, size=9, color=INK, line_spacing=1.2)

    # legend (bottom-left of chip band)
    legend_y = Inches(1.6)
    legend_items = [
        ("Deep", DEEP), ("Lab", LAB), ("Light", LIGHT), ("Meeting / ritual", MEETING)
    ]
    lx = SLIDE_W - MARGIN - Inches(5.2)
    for label, color in legend_items:
        sw = Inches(0.18)
        add_rrect(slide, lx, legend_y + Inches(0.04),
                  sw, Inches(0.18), fill=color, line_color=None)
        add_textbox(slide, lx + Inches(0.24), legend_y,
                    Inches(1.3), Inches(0.25),
                    label, size=9, color=MUTED, line_spacing=1.0)
        lx += Inches(1.3)


# ---------- slide 2: What you see each morning (dashboard + standup) ----------

def build_slide_2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(slide, 2, "What you see each morning",
               eyebrow="Daily touchpoint")

    add_textbox(slide, MARGIN, Inches(1.42),
                Inches(12.33), Inches(0.5),
                "Two surfaces you'd open first thing: the dashboard (in a browser) and the standup (in Claude chat). Mockups below. Tell us what's missing or extra.",
                size=12, color=MUTED, line_spacing=1.2)

    # ----- left: dashboard mockup -----
    dx = MARGIN
    dy = Inches(2.0)
    dw = Inches(6.2)
    dh = Inches(4.95)

    # card + fake browser chrome
    add_rrect(slide, dx, dy, dw, dh, fill=CARD_BG,
              line_color=LINE, line_width=Pt(0.75))
    add_rect(slide, dx, dy, dw, Inches(0.3),
             fill=tint(LINE, 0.5), line_color=None)
    for i, c in enumerate([RGBColor(0xff, 0x5f, 0x57),
                           RGBColor(0xfe, 0xbc, 0x2e),
                           RGBColor(0x28, 0xc8, 0x40)]):
        add_dot(slide, dx + Inches(0.15 + 0.17 * i), dy + Inches(0.15),
                Inches(0.1), c)
    add_rrect(slide, dx + Inches(0.7), dy + Inches(0.07),
              dw - Inches(0.9), Inches(0.16),
              fill=WHITE, line_color=LINE, line_width=Pt(0.4))
    add_textbox(slide, dx + Inches(0.8), dy + Inches(0.07),
                dw - Inches(1.1), Inches(0.16),
                "research-pm.github.io/dashboard",
                size=8, color=MUTED, line_spacing=1.0)

    cx = dx + Inches(0.25)
    cw = dw - Inches(0.5)
    cy = dy + Inches(0.45)

    # title row
    add_textbox(slide, cx, cy, cw, Inches(0.3),
                "research-pm — Tuesday, Apr 21",
                size=13, bold=True, color=NAVY, line_spacing=1.0)
    cy += Inches(0.35)

    # KPI row
    kpis = [
        ("ACTIVE PROJECTS", "3 / 3", "at WIP cap"),
        ("WRITING STREAK",  "12 days", "block held"),
        ("OVERDUE",         "1",      "from Monday"),
    ]
    kpi_w = (cw - Inches(0.3)) / 3
    for i, (label, value, sub) in enumerate(kpis):
        kx = cx + (kpi_w + Inches(0.15)) * i
        add_rrect(slide, kx, cy, kpi_w, Inches(0.65),
                  fill=tint(TEAL, 0.92), line_color=None)
        add_textbox(slide, kx + Inches(0.1), cy + Inches(0.05),
                    kpi_w - Inches(0.2), Inches(0.15),
                    label, size=7, bold=True, color=TEAL, line_spacing=1.0)
        add_textbox(slide, kx + Inches(0.1), cy + Inches(0.17),
                    kpi_w - Inches(0.2), Inches(0.3),
                    value, size=15, bold=True, color=NAVY, line_spacing=1.0)
        add_textbox(slide, kx + Inches(0.1), cy + Inches(0.44),
                    kpi_w - Inches(0.2), Inches(0.2),
                    sub, size=8, color=MUTED, line_spacing=1.0)
    cy += Inches(0.8)

    # Projects section
    add_textbox(slide, cx, cy, cw, Inches(0.22),
                "ACTIVE PROJECTS", size=9, bold=True, color=TEAL, line_spacing=1.0)
    cy += Inches(0.25)
    projects = [
        (GREEN, "TNF paper",       "last · methods draft (2h)",  "next · figures v3"),
        (AMBER, "Organoid model",  "last · media v2 protocol",   "next · cell type B stable"),
        (GREEN, "R01 renewal",     "last · Aim 3 outline",       "next · preliminary data"),
    ]
    for status, name, last, nxt in projects:
        add_line(slide, cx, cy + Inches(0.52),
                 cx + cw, cy + Inches(0.52),
                 color=tint(LINE, 0.5), width_pt=0.5)
        add_dot(slide, cx + Inches(0.1), cy + Inches(0.16),
                Inches(0.12), status)
        add_textbox(slide, cx + Inches(0.25), cy + Inches(0.03),
                    Inches(2.5), Inches(0.22),
                    name, size=11, bold=True, color=NAVY, line_spacing=1.0)
        add_textbox(slide, cx + Inches(0.25), cy + Inches(0.26),
                    cw - Inches(0.25), Inches(0.2),
                    f"{last}    ·    {nxt}",
                    size=8, color=MUTED, line_spacing=1.0)
        cy += Inches(0.53)
    cy += Inches(0.05)

    # Quarterly goals
    add_textbox(slide, cx, cy, cw, Inches(0.22),
                "QUARTERLY GOALS", size=9, bold=True, color=TEAL, line_spacing=1.0)
    cy += Inches(0.25)
    goals_q = [
        ("TNF paper  ·  Q2 draft  (OKR)",      0.67, AMBER),
        ("Organoid  ·  Q2 baseline  (narr.)",  0.40, AMBER),
    ]
    for label, pct, color in goals_q:
        add_textbox(slide, cx, cy, Inches(3.2), Inches(0.2),
                    label, size=9, color=INK, line_spacing=1.0)
        bar_x = cx + Inches(3.2)
        bar_w = cw - Inches(3.6)
        add_rect(slide, bar_x, cy + Inches(0.06),
                 bar_w, Inches(0.1),
                 fill=tint(LINE, 0.3), line_color=None)
        add_rect(slide, bar_x, cy + Inches(0.06),
                 Emu(int(bar_w * pct)), Inches(0.1),
                 fill=color, line_color=None)
        add_textbox(slide, cx + cw - Inches(0.35), cy,
                    Inches(0.35), Inches(0.2),
                    f"{int(pct*100)}%", size=8, bold=True, color=color,
                    align=PP_ALIGN.RIGHT, line_spacing=1.0)
        cy += Inches(0.23)

    # ----- right: standup (Claude chat) mockup -----
    sx = dx + dw + Inches(0.2)
    sy = dy
    sw = SLIDE_W - MARGIN - sx
    sh = dh

    add_rrect(slide, sx, sy, sw, sh,
              fill=tint(NAVY, 0.97), line_color=LINE, line_width=Pt(0.75))
    add_rect(slide, sx, sy, sw, Inches(0.35),
             fill=NAVY, line_color=None)
    add_dot(slide, sx + Inches(0.2), sy + Inches(0.175),
            Inches(0.16), TEAL)
    add_textbox(slide, sx + Inches(0.4), sy + Inches(0.08),
                sw - Inches(0.5), Inches(0.22),
                "Claude  ·  morning standup",
                size=10, bold=True, color=WHITE, line_spacing=1.0)
    add_textbox(slide, sx + sw - Inches(1.5), sy + Inches(0.08),
                Inches(1.3), Inches(0.22),
                "9:02 am", size=8, color=tint(TEAL, 0.5),
                align=PP_ALIGN.RIGHT, line_spacing=1.0)

    bx = sx + Inches(0.25)
    bw = sw - Inches(0.5)
    by = sy + Inches(0.5)

    add_textbox(slide, bx, by, bw, Inches(0.25),
                "Good morning — standup for Tuesday, Apr 21.",
                size=11, color=INK, line_spacing=1.0)
    by += Inches(0.35)

    # TODAY'S 3
    add_textbox(slide, bx, by, bw, Inches(0.22),
                "TODAY'S 3", size=9, bold=True, color=TEAL, line_spacing=1.0)
    by += Inches(0.24)
    todays = [
        ("09:00–11:00", DEEP,    "Methods revise",
         "TNF paper  ›  Methods drafted  ›  Q2 draft complete"),
        ("14:00–16:00", LAB,     "Media protocol v2",
         "Organoid  ›  Media v2  ›  Q2 baseline"),
        ("16:15–17:00", MEETING, "1:1 with Sara",
         "People  ›  Sara owes 2  ›  4-week cadence"),
    ]
    for tm, color, task, chain in todays:
        add_rect(slide, bx, by + Inches(0.03),
                 Inches(0.1), Inches(0.22), fill=color, line_color=None)
        add_textbox(slide, bx + Inches(0.17), by,
                    Inches(1.25), Inches(0.22),
                    tm, size=8, bold=True, color=MUTED, line_spacing=1.0,
                    font="Consolas")
        add_textbox(slide, bx + Inches(1.5), by,
                    bw - Inches(1.5), Inches(0.22),
                    task, size=10, bold=True, color=NAVY, line_spacing=1.0)
        add_textbox(slide, bx + Inches(0.17), by + Inches(0.22),
                    bw - Inches(0.17), Inches(0.2),
                    chain, size=7.5, color=TEAL, line_spacing=1.0)
        by += Inches(0.46)

    by += Inches(0.05)

    # Overdue
    add_textbox(slide, bx, by, bw, Inches(0.22),
                "OVERDUE", size=9, bold=True, color=RED, line_spacing=1.0)
    by += Inches(0.22)
    add_textbox(slide, bx + Inches(0.1), by,
                bw - Inches(0.1), Inches(0.22),
                "Figure 4 v1 (TNF) — due Monday. Re-slot, or push?",
                size=9, color=INK, line_spacing=1.0)
    by += Inches(0.3)

    # Calendar shape
    add_textbox(slide, bx, by, bw, Inches(0.22),
                "CALENDAR SHAPE", size=9, bold=True, color=TEAL, line_spacing=1.0)
    by += Inches(0.22)
    add_textbox(slide, bx + Inches(0.1), by,
                bw - Inches(0.1), Inches(0.22),
                "3h Deep morning  ·  2h Lab afternoon  ·  1 meeting",
                size=9, color=INK, line_spacing=1.0)
    by += Inches(0.3)

    # Heads up
    add_textbox(slide, bx, by, bw, Inches(0.22),
                "HEADS UP", size=9, bold=True, color=AMBER, line_spacing=1.0)
    by += Inches(0.22)
    add_textbox(slide, bx + Inches(0.1), by,
                bw - Inches(0.1), Inches(0.4),
                "Writing block skipped twice this week. Organoid log silent 4 days.",
                size=9, color=INK, line_spacing=1.2)


# ---------- slide 3: Goal cascade ----------

def build_slide_3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(slide, 3, "How goals connect",
               eyebrow="Goal structure")

    add_textbox(slide, MARGIN, Inches(1.42),
                Inches(12.33), Inches(0.5),
                "Year goals break into quarters, then projects, then weeks, then days. Quarterly goals come in two styles: OKR (Objective + measurable Key Results) for tight-clock work; narrative (status + paragraph) for exploratory work.",
                size=12, color=MUTED, line_spacing=1.2)

    # level labels on the left
    level_labels = [
        (Inches(2.05), "YEAR"),
        (Inches(2.95), "QUARTER"),
        (Inches(4.1),  "PROJECT"),
        (Inches(5.4),  "WEEK"),
        (Inches(6.25), "TODAY"),
    ]
    for y, label in level_labels:
        add_textbox(slide, MARGIN - Inches(0.1), y,
                    Inches(1.0), Inches(0.25),
                    label, size=8, bold=True, color=TEAL,
                    line_spacing=1.0)

    # two column centers
    left_cx = Inches(4.3)
    right_cx = Inches(9.0)
    box_w = Inches(3.8)

    # -------- row 1: year goals (two) --------
    y = Inches(2.0)
    h = Inches(0.7)

    # left year
    lx = left_cx - box_w/2
    add_rrect(slide, lx, y, box_w, h, fill=NAVY, line_color=None)
    add_textbox(slide, lx + Inches(0.2), y + Inches(0.08),
                box_w - Inches(0.4), Inches(0.25),
                "YEAR GOAL · OKR", size=8, bold=True, color=tint(TEAL, 0.4),
                line_spacing=1.0)
    add_textbox(slide, lx + Inches(0.2), y + Inches(0.32),
                box_w - Inches(0.4), Inches(0.35),
                "Publish TNF paper (Cell Reports, Oct 2026)",
                size=13, bold=True, color=WHITE, line_spacing=1.1)

    # right year
    rx = right_cx - box_w/2
    add_rrect(slide, rx, y, box_w, h, fill=NAVY, line_color=None)
    add_textbox(slide, rx + Inches(0.2), y + Inches(0.08),
                box_w - Inches(0.4), Inches(0.25),
                "YEAR GOAL · NARRATIVE", size=8, bold=True, color=tint(TEAL, 0.4),
                line_spacing=1.0)
    add_textbox(slide, rx + Inches(0.2), y + Inches(0.32),
                box_w - Inches(0.4), Inches(0.35),
                "Establish organoid model — 2 cell types, protocol",
                size=13, bold=True, color=WHITE, line_spacing=1.1)

    # connectors from year → quarter
    add_arrow(slide, left_cx, y + h, left_cx, Inches(2.9), color=MUTED)
    add_arrow(slide, right_cx, y + h, right_cx, Inches(2.9), color=MUTED)

    # -------- row 2: quarterly (OKR on left, narrative on right) --------
    y = Inches(2.9)
    h = Inches(1.05)

    # left: OKR card with KRs
    add_rrect(slide, lx, y, box_w, h, fill=CARD_BG,
              line_color=LINE, line_width=Pt(0.75))
    add_textbox(slide, lx + Inches(0.2), y + Inches(0.08),
                box_w - Inches(0.4), Inches(0.3),
                "Q2 — Draft complete  (OKR)", size=12, bold=True, color=NAVY,
                line_spacing=1.0)

    # 3 KRs as mini progress bars
    kr_defs = [
        ("Key Result 1 — Methods drafted",  1.00, GREEN),
        ("Key Result 2 — Figures v3",       0.67, AMBER),
        ("Key Result 3 — Co-author review", 0.00, MUTED),
    ]
    for i, (kr_label, pct, color) in enumerate(kr_defs):
        ky = y + Inches(0.4 + i * 0.21)
        add_textbox(slide, lx + Inches(0.2), ky,
                    Inches(1.8), Inches(0.2),
                    kr_label, size=9, color=INK, line_spacing=1.0)
        # bar bg
        bar_x = lx + Inches(2.1)
        bar_w = Inches(1.4)
        add_rect(slide, bar_x, ky + Inches(0.05),
                 bar_w, Inches(0.1),
                 fill=tint(LINE, 0.3), line_color=None)
        if pct > 0:
            add_rect(slide, bar_x, ky + Inches(0.05),
                     Emu(int(bar_w * pct)), Inches(0.1),
                     fill=color, line_color=None)
        add_textbox(slide, lx + box_w - Inches(0.5), ky,
                    Inches(0.35), Inches(0.2),
                    f"{int(pct*100)}%", size=9, bold=True, color=color,
                    align=PP_ALIGN.RIGHT, line_spacing=1.0)

    # right: narrative card
    add_rrect(slide, rx, y, box_w, h, fill=CARD_BG,
              line_color=LINE, line_width=Pt(0.75))
    add_textbox(slide, rx + Inches(0.2), y + Inches(0.08),
                Inches(3), Inches(0.3),
                "Q2 — Baseline working  (narrative)", size=12, bold=True, color=NAVY,
                line_spacing=1.0)
    # traffic light
    add_dot(slide, rx + box_w - Inches(0.35), y + Inches(0.22),
            Inches(0.18), AMBER)
    add_textbox(slide, rx + Inches(0.2), y + Inches(0.42),
                box_w - Inches(0.4), Inches(0.55),
                "Media reliably reproducible. 2nd cell type has recurring contamination — root cause traced to shared incubator. Next: isolate. Status: on track.",
                size=9, color=INK, line_spacing=1.25)

    # connectors quarter → project
    add_arrow(slide, left_cx, y + h, left_cx, Inches(4.05), color=MUTED)
    add_arrow(slide, right_cx, y + h, right_cx, Inches(4.05), color=MUTED)

    # -------- row 3: project checkpoints --------
    y = Inches(4.05)
    h = Inches(0.85)

    for cx, items in [
        (left_cx, [("Methods drafted", GREEN, 1.0),
                   ("Results drafted", AMBER, 0.5),
                   ("Figures v3", MUTED, 0.0)]),
        (right_cx, [("Media protocol v2", GREEN, 1.0),
                    ("Cell type B stable", AMBER, 0.4),
                    ("Doubling time char.", MUTED, 0.0)]),
    ]:
        px = cx - box_w/2
        add_rrect(slide, px, y, box_w, h, fill=CARD_BG,
                  line_color=LINE, line_width=Pt(0.75))
        add_textbox(slide, px + Inches(0.2), y + Inches(0.08),
                    box_w - Inches(0.4), Inches(0.22),
                    "CHECKPOINTS  ·  each has a Definition of Done",
                    size=8, bold=True, color=TEAL, line_spacing=1.0)
        # 3 mini-rows
        for i, (label, color, pct) in enumerate(items):
            ly = y + Inches(0.32 + i*0.17)
            add_dot(slide, px + Inches(0.3), ly + Inches(0.08),
                    Inches(0.1), color)
            add_textbox(slide, px + Inches(0.45), ly,
                        box_w - Inches(0.6), Inches(0.18),
                        label, size=9, color=INK, line_spacing=1.0)

    # connectors → weekly (converge)
    add_arrow(slide, left_cx, y + h, Inches(6.667) - Inches(0.8), Inches(5.25),
              color=MUTED)
    add_arrow(slide, right_cx, y + h, Inches(6.667) + Inches(0.8), Inches(5.25),
              color=MUTED)

    # -------- row 4: weekly top-5 --------
    y = Inches(5.25)
    h = Inches(0.75)
    w = Inches(6.5)
    wx = Inches(6.667) - w/2
    add_rrect(slide, wx, y, w, h, fill=TEAL, line_color=None)
    add_textbox(slide, wx + Inches(0.2), y + Inches(0.08),
                w - Inches(0.4), Inches(0.25),
                "WEEKLY TOP-5  ·  max 1–2 per project", size=9, bold=True,
                color=tint(TEAL, 0.7), line_spacing=1.0)
    weekly = [
        "1 · Methods revise (TNF)",
        "2 · Figure 4 v1 (TNF)",
        "3 · Media protocol v2 (Organoid)",
        "4 · 1:1 Sara (postdoc)",
        "5 · R01 Aim 3 data outline",
    ]
    for i, w_label in enumerate(weekly):
        col = i % 3
        row = i // 3
        col_w = (w - Inches(0.4)) / 3
        add_textbox(slide, wx + Inches(0.2) + col_w * col,
                    y + Inches(0.35 + row*0.22),
                    col_w, Inches(0.22),
                    w_label, size=10, bold=True, color=WHITE, line_spacing=1.0)

    # connector → today
    add_arrow(slide, Inches(6.667), y + h, Inches(6.667), Inches(6.18),
              color=MUTED)

    # -------- row 5: today's 3 --------
    y = Inches(6.18)
    h = Inches(0.65)
    tw = Inches(5.5)
    tx = Inches(6.667) - tw/2
    add_rrect(slide, tx, y, tw, h, fill=NAVY, line_color=None)
    add_textbox(slide, tx + Inches(0.2), y + Inches(0.08),
                tw - Inches(0.4), Inches(0.22),
                "TODAY'S 3  ·  slotted into calendar blocks", size=9, bold=True,
                color=tint(TEAL, 0.4), line_spacing=1.0)
    today_items = [
        ("09:00–11:00  Methods revise",   DEEP),
        ("14:00–16:00  Media protocol",   LAB),
        ("16:15–17:00  1:1 Sara",         MEETING),
    ]
    col_w = (tw - Inches(0.4)) / 3
    for i, (label, color) in enumerate(today_items):
        cx = tx + Inches(0.2) + col_w * i
        add_rect(slide, cx, y + Inches(0.32), Inches(0.12), Inches(0.2),
                 fill=color, line_color=None)
        add_textbox(slide, cx + Inches(0.18), y + Inches(0.3),
                    col_w - Inches(0.2), Inches(0.3),
                    label, size=10, color=WHITE, line_spacing=1.0)


# ---------- slide 4: multi-level Gantt ----------

def build_slide_4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(slide, 4, "One data source, two views",
               eyebrow="Visualise")

    add_textbox(slide, MARGIN, Inches(1.42),
                Inches(12.33), Inches(0.5),
                "All views come from the same markdown files. The Gantt (top) shows a project across months. The Plan Board (bottom) shows this week's calendar blocks. Move something in one view, the other updates.",
                size=12, color=MUTED, line_spacing=1.2)

    # ===== top: project Gantt with swim lanes =====
    gantt_x = MARGIN
    gantt_y = Inches(1.95)
    gantt_w = Inches(12.33)
    gantt_h = Inches(2.7)

    # label strip
    add_textbox(slide, gantt_x, gantt_y - Inches(0.05),
                Inches(5), Inches(0.25),
                "PROJECT GANTT · TNF paper", size=9, bold=True, color=TEAL,
                line_spacing=1.0)

    # card shell
    add_rrect(slide, gantt_x, gantt_y + Inches(0.18),
              gantt_w, gantt_h - Inches(0.18),
              fill=CARD_BG, line_color=LINE, line_width=Pt(0.75))

    # lane area
    lane_label_w = Inches(1.2)
    timeline_x0 = gantt_x + lane_label_w + Inches(0.2)
    timeline_w = gantt_w - lane_label_w - Inches(0.4)
    month_w = timeline_w / 3  # Apr, May, Jun

    # month headers
    months = ["APR", "MAY", "JUN"]
    for i, m in enumerate(months):
        mx = timeline_x0 + month_w * i
        # vertical gridline
        add_line(slide, mx, gantt_y + Inches(0.55),
                 mx, gantt_y + gantt_h - Inches(0.1),
                 color=tint(LINE, 0.3), width_pt=0.5)
        add_textbox(slide, mx, gantt_y + Inches(0.3),
                    month_w, Inches(0.25),
                    m, size=9, bold=True, color=MUTED,
                    align=PP_ALIGN.LEFT, line_spacing=1.0)
    # final gridline
    add_line(slide, timeline_x0 + timeline_w,
             gantt_y + Inches(0.55),
             timeline_x0 + timeline_w,
             gantt_y + gantt_h - Inches(0.1),
             color=tint(LINE, 0.3), width_pt=0.5)

    # swim lanes
    lanes = [
        ("Writing",  DEEP, [
            ("Intro", 0, 14, True, False),        # done
            ("Methods", 14, 30, False, True),     # active
            ("Results", 30, 45, False, False),
            ("Figures v3", 45, 62, False, False, True),  # crit
        ]),
        ("Lab",      LAB, [
            ("Dose-response", 20, 45, False, False),
            ("Replicate", 45, 60, False, False),
        ]),
        ("Collab",   MEETING, [
            ("Co-author review", 76, 0, False, False, False, True),  # milestone
        ]),
    ]
    # Total timeline days: Apr 1 - Jun 30 = 91 days
    total_days = 91
    lane_h = (gantt_h - Inches(0.65)) / len(lanes)
    lane_y0 = gantt_y + Inches(0.55)

    methods_bar_rect = None  # save for connector

    for li, lane in enumerate(lanes):
        lane_name, lane_color, bars = lane[0], lane[1], lane[2]
        ly = lane_y0 + lane_h * li
        # lane label
        add_textbox(slide, gantt_x + Inches(0.15),
                    ly + lane_h/2 - Inches(0.12),
                    lane_label_w, Inches(0.25),
                    lane_name, size=10, bold=True, color=INK,
                    line_spacing=1.0)
        # lane divider
        if li > 0:
            add_line(slide, gantt_x + Inches(0.05), ly,
                     gantt_x + gantt_w - Inches(0.05), ly,
                     color=tint(LINE, 0.4), width_pt=0.5)
        # bars
        for bar in bars:
            name = bar[0]
            start_day = bar[1]
            span_day = bar[2]
            done = bar[3]
            active = bar[4] if len(bar) > 4 else False
            crit = bar[5] if len(bar) > 5 else False
            milestone = bar[6] if len(bar) > 6 else False

            bx = timeline_x0 + (timeline_w * start_day / total_days)
            by = ly + lane_h/2 - Inches(0.18)
            bh = Inches(0.36)

            if milestone:
                # diamond
                dm = slide.shapes.add_shape(MSO_SHAPE.DIAMOND,
                                            bx - Inches(0.18), by,
                                            Inches(0.36), bh)
                set_solid_fill(dm, lane_color)
                set_line(dm, None)
                add_textbox(slide, bx + Inches(0.25), by + Inches(0.06),
                            Inches(1.8), Inches(0.25),
                            name, size=9, color=INK, line_spacing=1.0)
                continue

            bw = Emu(int(timeline_w * span_day / total_days))

            # fill logic
            if done:
                fill = tint(lane_color, 0.55)
                text_color = INK
            elif crit:
                fill = lane_color
                text_color = WHITE
            elif active:
                fill = lane_color
                text_color = WHITE
            else:
                fill = tint(lane_color, 0.25)
                text_color = INK

            rect = add_rrect(slide, bx, by, bw, bh,
                             fill=fill, line_color=None)

            if active and name == "Methods":
                methods_bar_rect = (bx, by, bw, bh)

            # label inside or to the right
            if bw > Inches(0.7):
                add_textbox(slide, bx + Inches(0.08), by + Inches(0.06),
                            bw - Inches(0.16), Inches(0.25),
                            name, size=9, bold=True, color=text_color,
                            line_spacing=1.0)
            else:
                add_textbox(slide, bx + bw + Inches(0.05), by + Inches(0.06),
                            Inches(1.5), Inches(0.25),
                            name, size=9, color=INK, line_spacing=1.0)

    # ===== bottom: Plan Board week tile =====
    board_y = Inches(4.95)
    board_h = Inches(1.95)

    add_textbox(slide, gantt_x, board_y - Inches(0.05),
                Inches(6), Inches(0.25),
                "PLAN BOARD · week of Apr 20", size=9, bold=True, color=TEAL,
                line_spacing=1.0)

    board_cx = gantt_x
    board_w = gantt_w
    add_rrect(slide, board_cx, board_y + Inches(0.18),
              board_w, board_h - Inches(0.18),
              fill=CARD_BG, line_color=LINE, line_width=Pt(0.75))

    # time rail on left
    time_rail_x = board_cx + Inches(0.15)
    board_day_x0 = board_cx + Inches(0.85)
    board_day_w = (board_w - Inches(1.0)) / 5
    day_labels = ["MON 20", "TUE 21", "WED 22", "THU 23", "FRI 24"]
    for i, d in enumerate(day_labels):
        dx = board_day_x0 + board_day_w * i
        add_textbox(slide, dx, board_y + Inches(0.3),
                    board_day_w, Inches(0.22),
                    d, size=9, bold=True, color=MUTED,
                    align=PP_ALIGN.CENTER, line_spacing=1.0)

    # time labels + gridlines
    time_labels = ["09", "11", "13", "15", "17"]
    top_time = board_y + Inches(0.55)
    time_h = (board_h - Inches(0.7))
    for i, t in enumerate(time_labels):
        ty = top_time + time_h * (i / (len(time_labels)))
        add_textbox(slide, time_rail_x, ty,
                    Inches(0.6), Inches(0.2),
                    t, size=8, color=MUTED, line_spacing=1.0)
        add_line(slide, board_day_x0, ty + Inches(0.1),
                 board_day_x0 + board_day_w * 5, ty + Inches(0.1),
                 color=tint(LINE, 0.4), width_pt=0.5)

    # day-column dividers
    for i in range(6):
        lx = board_day_x0 + board_day_w * i
        add_line(slide, lx, top_time,
                 lx, top_time + time_h,
                 color=tint(LINE, 0.3), width_pt=0.5)

    # Plan Board blocks (day_idx, top_pct, height_pct, color, label)
    # top_pct/height_pct are fraction of time_h
    blocks = [
        # Mon
        (0, 0.00, 0.30, DEEP,    "Methods"),
        (0, 0.50, 0.15, LIGHT,   "Admin"),
        # Tue — this is the Methods block we connect to
        (1, 0.00, 0.30, DEEP,    "Methods"),
        (1, 0.45, 0.25, LAB,     "Dose-resp"),
        (1, 0.80, 0.12, MEETING, "Ravi"),
        # Wed
        (2, 0.00, 0.32, DEEP,    "Results"),
        (2, 0.45, 0.30, LAB,     "Media"),
        (2, 0.80, 0.15, MEETING, "Lab mtg"),
        # Thu
        (3, 0.00, 0.28, DEEP,    "Revise"),
        (3, 0.45, 0.22, LAB,     "Assay"),
        (3, 0.80, 0.12, MEETING, "Sara"),
        # Fri
        (4, 0.00, 0.28, DEEP,    "Polish"),
        (4, 0.50, 0.12, LAB,     "Passage"),
        (4, 0.80, 0.12, MEETING, "/digest"),
    ]

    tue_methods_rect = None
    for day_idx, top_pct, h_pct, color, label in blocks:
        bx = board_day_x0 + board_day_w * day_idx + Inches(0.05)
        by = top_time + time_h * top_pct
        bw = board_day_w - Inches(0.1)
        bh = Emu(int(time_h * h_pct))
        add_rrect(slide, bx, by, bw, bh, fill=color, line_color=None)
        add_textbox(slide, bx + Inches(0.06), by + Inches(0.05),
                    bw - Inches(0.12), bh - Inches(0.1),
                    label, size=8, bold=True, color=WHITE, line_spacing=1.0)
        if day_idx == 1 and label == "Methods":
            tue_methods_rect = (bx, by, bw, bh)

    # ===== dashed connector: Methods Gantt bar → Tuesday Methods block =====
    if methods_bar_rect and tue_methods_rect:
        # from bottom-center of Gantt methods bar
        gx = methods_bar_rect[0] + methods_bar_rect[2] / 2
        gy = methods_bar_rect[1] + methods_bar_rect[3]
        # to top-center of Tue methods block
        tx = tue_methods_rect[0] + tue_methods_rect[2] / 2
        ty = tue_methods_rect[1]
        add_dashed_line(slide, gx, gy, tx, ty, color=TEAL, width_pt=1.5)
        # small label on the connector
        mid_y = (gy + ty) / 2
        add_textbox(slide, gx - Inches(0.15), mid_y - Inches(0.12),
                    Inches(2.5), Inches(0.25),
                    "same markdown", size=8, bold=True, color=TEAL,
                    line_spacing=1.0)


# ---------- slide 5: Move → cascade ----------

def build_slide_5(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(slide, 5, "Moving things around",
               eyebrow="Modify")

    add_textbox(slide, MARGIN, Inches(1.42),
                Inches(12.33), Inches(0.5),
                "Drag the bar on the Gantt, or tell Claude in chat. Either way, every downstream block shifts — and you see the impact before anything saves. (\"Cascade\" = the knock-on effect.)",
                size=12, color=MUTED, line_spacing=1.2)

    # ===== left panel: BEFORE =====
    panel_y = Inches(2.0)
    panel_h = Inches(2.3)
    panel_w = Inches(5.9)
    before_x = MARGIN

    add_textbox(slide, before_x, panel_y - Inches(0.05),
                Inches(3), Inches(0.25),
                "BEFORE", size=9, bold=True, color=MUTED,
                line_spacing=1.0)
    add_rrect(slide, before_x, panel_y + Inches(0.18),
              panel_w, panel_h - Inches(0.18),
              fill=CARD_BG, line_color=LINE, line_width=Pt(0.75))

    # mini gantt: 3 bars stacked
    # timeline: Apr 15 - Jun 15 = 62 days
    def draw_mini_gantt(slide, x, y, w, h, bars, highlight_ghost_bars=None):
        """bars: list of (lane_name, start, span, label, color, status)
        status: 'done', 'active', 'ghost', 'shifted', 'plain'
        highlight_ghost_bars: separate list of ghost (original) positions
        """
        total = 62
        lane_label_w = Inches(0.7)
        tx0 = x + lane_label_w + Inches(0.15)
        tw = w - lane_label_w - Inches(0.3)
        # month markers
        months = [("APR 15", 0), ("MAY 1", 16), ("MAY 15", 30), ("JUN 1", 47), ("JUN 15", 61)]
        for m, d in months:
            mx = tx0 + tw * d / total
            add_line(slide, mx, y + Inches(0.45),
                     mx, y + h - Inches(0.1),
                     color=tint(LINE, 0.35), width_pt=0.4)
        for m, d in [("APR", 0), ("MAY", 16), ("JUN", 47)]:
            mx = tx0 + tw * d / total
            add_textbox(slide, mx + Inches(0.05), y + Inches(0.2),
                        Inches(0.7), Inches(0.2),
                        m, size=8, bold=True, color=MUTED, line_spacing=1.0)

        unique_lanes = []
        for b in bars:
            if b[0] not in unique_lanes:
                unique_lanes.append(b[0])
        lane_h = (h - Inches(0.5)) / len(unique_lanes)
        lane_y0 = y + Inches(0.45)

        # ghost bars first
        if highlight_ghost_bars:
            for lane_name, start, span, label, color in highlight_ghost_bars:
                li = unique_lanes.index(lane_name)
                ly = lane_y0 + lane_h * li + lane_h/2 - Inches(0.14)
                bx = tx0 + tw * start / total
                bw = Emu(int(tw * span / total))
                ghost = add_rect(slide, bx, ly, bw, Inches(0.28),
                                 fill=None, line_color=tint(color, 0.6),
                                 line_width=Pt(1.0))
                # dashed outline
                set_line_dashed(ghost, tint(color, 0.5), width_pt=1.0)

        for bar in bars:
            lane_name, start, span, label, color, status = bar
            li = unique_lanes.index(lane_name)
            # lane label
            add_textbox(slide, x + Inches(0.1),
                        lane_y0 + lane_h * li + lane_h/2 - Inches(0.1),
                        lane_label_w, Inches(0.2),
                        lane_name, size=8, bold=True, color=INK,
                        line_spacing=1.0)

            ly = lane_y0 + lane_h * li + lane_h/2 - Inches(0.14)
            bx = tx0 + tw * start / total
            bw = Emu(int(tw * span / total))

            if status == "done":
                fill = tint(color, 0.55)
                tcolor = INK
            elif status == "active":
                fill = color
                tcolor = WHITE
            elif status == "shifted":
                fill = color
                tcolor = WHITE
            elif status == "crit":
                fill = color
                tcolor = WHITE
            else:
                fill = tint(color, 0.3)
                tcolor = INK

            rect = add_rrect(slide, bx, ly, bw, Inches(0.28),
                             fill=fill, line_color=None)
            if bw > Inches(0.6):
                add_textbox(slide, bx + Inches(0.06), ly + Inches(0.02),
                            bw - Inches(0.12), Inches(0.24),
                            label, size=8, bold=True, color=tcolor,
                            line_spacing=1.0)

    draw_mini_gantt(
        slide,
        before_x + Inches(0.15), panel_y + Inches(0.3),
        panel_w - Inches(0.3), panel_h - Inches(0.4),
        bars=[
            ("Writing", 0, 15, "Methods",    DEEP,    "active"),
            ("Writing", 15, 15, "Results",   DEEP,    "plain"),
            ("Writing", 30, 17, "Figures v3", DEEP,   "crit"),
            ("Collab",  61, 1,  "Review",    MEETING, "plain"),
        ],
    )

    # ===== right panel: AFTER =====
    after_x = SLIDE_W - MARGIN - panel_w
    add_textbox(slide, after_x, panel_y - Inches(0.05),
                Inches(3), Inches(0.25),
                "AFTER", size=9, bold=True, color=TEAL,
                line_spacing=1.0)
    add_rrect(slide, after_x, panel_y + Inches(0.18),
              panel_w, panel_h - Inches(0.18),
              fill=CARD_BG, line_color=LINE, line_width=Pt(0.75))

    draw_mini_gantt(
        slide,
        after_x + Inches(0.15), panel_y + Inches(0.3),
        panel_w - Inches(0.3), panel_h - Inches(0.4),
        bars=[
            # methods shifted +7
            ("Writing", 7, 15, "Methods",    DEEP,    "shifted"),
            ("Writing", 22, 15, "Results",   DEEP,    "shifted"),
            ("Writing", 37, 17, "Figures v3", DEEP,   "crit"),
            ("Collab",  61, 1,  "Review",    MEETING, "plain"),
        ],
        highlight_ghost_bars=[
            ("Writing", 0, 15, "Methods", DEEP),
            ("Writing", 15, 15, "Results", DEEP),
            ("Writing", 30, 17, "Figures v3", DEEP),
        ],
    )

    # ===== center "Drag OR Say" callout between panels =====
    # Actually with the wider panels there's no center gap. Put callout below or above.
    # Put a slim center column between them? Let me put as bottom-overlay badge.

    # ===== middle strip: two input modes =====
    ms_y = Inches(4.45)
    ms_h = Inches(0.8)

    # left: Drag
    add_rrect(slide, MARGIN, ms_y, panel_w, ms_h,
              fill=tint(TEAL, 0.9), line_color=TEAL, line_width=Pt(0.75))
    add_textbox(slide, MARGIN + Inches(0.25), ms_y + Inches(0.1),
                Inches(1.5), Inches(0.25),
                "DRAG", size=9, bold=True, color=TEAL, line_spacing=1.0)
    add_textbox(slide, MARGIN + Inches(0.25), ms_y + Inches(0.35),
                panel_w - Inches(0.5), Inches(0.45),
                "Grab the Methods bar, drop it a week later on the Gantt or Plan Board.",
                size=11, color=INK, line_spacing=1.2)

    # center "OR"
    or_x = MARGIN + panel_w + Inches(0.1)
    or_w = Inches(0.5)
    add_dot(slide, or_x + or_w/2, ms_y + ms_h/2, Inches(0.5), NAVY)
    add_textbox(slide, or_x, ms_y + ms_h/2 - Inches(0.13),
                or_w, Inches(0.3),
                "OR", size=11, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, line_spacing=1.0)

    # right: Say
    say_x = SLIDE_W - MARGIN - panel_w
    add_rrect(slide, say_x, ms_y, panel_w, ms_h,
              fill=tint(TEAL, 0.9), line_color=TEAL, line_width=Pt(0.75))
    add_textbox(slide, say_x + Inches(0.25), ms_y + Inches(0.1),
                Inches(1.5), Inches(0.25),
                "SAY", size=9, bold=True, color=TEAL, line_spacing=1.0)
    add_textbox(slide, say_x + Inches(0.25), ms_y + Inches(0.35),
                panel_w - Inches(0.5), Inches(0.45),
                "\u201cslip TNF methods by a week\u201d in Claude chat.",
                size=11, color=INK, line_spacing=1.2, font="Consolas")

    # ===== impact panel =====
    imp_y = Inches(5.4)
    imp_h = Inches(1.05)
    imp_w = Inches(7.5)
    imp_x = MARGIN

    add_rrect(slide, imp_x, imp_y, imp_w, imp_h,
              fill=NAVY, line_color=None)
    add_textbox(slide, imp_x + Inches(0.3), imp_y + Inches(0.1),
                Inches(4), Inches(0.25),
                "IMPACT", size=9, bold=True, color=tint(TEAL, 0.4),
                line_spacing=1.0)

    # three rows of impact
    imp_items = [
        ("Q2 Objective",             "\u25cf \u2192 \u25cf", "on track \u2192 watch",   GREEN, AMBER),
        ("Year goal",                "\u25cf",               "on track",                 GREEN, None),
        ("Buffer on Figures v3",     "",                     "2 days remaining",         AMBER, None),
    ]

    for i, (label, glyph, text, c1, c2) in enumerate(imp_items):
        iy = imp_y + Inches(0.4) + Inches(0.22) * i
        add_textbox(slide, imp_x + Inches(0.3), iy,
                    Inches(2.2), Inches(0.2),
                    label, size=10, color=WHITE, line_spacing=1.0)
        # dots
        dx = imp_x + Inches(2.5)
        add_dot(slide, dx, iy + Inches(0.1), Inches(0.14), c1)
        if c2:
            add_textbox(slide, dx + Inches(0.12), iy - Inches(0.01),
                        Inches(0.5), Inches(0.22),
                        "\u2192", size=12, color=MUTED, line_spacing=1.0)
            add_dot(slide, dx + Inches(0.45), iy + Inches(0.1), Inches(0.14), c2)
        # text
        add_textbox(slide, dx + Inches(0.8), iy,
                    Inches(4.5), Inches(0.22),
                    text, size=10, bold=True, color=tint(TEAL, 0.6),
                    line_spacing=1.0)

    # ===== three action buttons =====
    btn_y = Inches(5.4)
    btn_h = Inches(1.05)
    btn_x = imp_x + imp_w + Inches(0.2)
    btn_w = SLIDE_W - MARGIN - btn_x
    btn_cell_h = (btn_h - Inches(0.1)) / 3

    buttons = [
        ("Compress downstream", "Tighten the next blocks to hit original date.", GREEN),
        ("Accept the slip",     "Push Figures v3 deadline; Co-author review stays.", AMBER),
        ("Cut scope",           "Drop 1 figure; stay on original timeline.", RED),
    ]
    for i, (label, body, color) in enumerate(buttons):
        y = btn_y + (btn_cell_h + Inches(0.05)) * i
        add_rrect(slide, btn_x, y, btn_w, btn_cell_h,
                  fill=CARD_BG, line_color=color, line_width=Pt(1.5))
        # color dot
        add_dot(slide, btn_x + Inches(0.25), y + btn_cell_h/2,
                Inches(0.2), color)
        add_textbox(slide, btn_x + Inches(0.55), y + Inches(0.04),
                    btn_w - Inches(0.7), Inches(0.28),
                    label, size=11, bold=True, color=NAVY, line_spacing=1.0)
        add_textbox(slide, btn_x + Inches(0.55), y + Inches(0.28),
                    btn_w - Inches(0.7), btn_cell_h - Inches(0.3),
                    body, size=9, color=MUTED, line_spacing=1.2)


# ---------- slide 6: Artifact gallery ----------

def _panel_shell(slide, x, y, w, h, eyebrow, title):
    add_rrect(slide, x, y, w, h, fill=CARD_BG,
              line_color=LINE, line_width=Pt(0.75))
    add_textbox(slide, x + Inches(0.2), y + Inches(0.1),
                w - Inches(0.4), Inches(0.2),
                eyebrow.upper(), size=8, bold=True, color=TEAL, line_spacing=1.0)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.28),
                w - Inches(0.4), Inches(0.3),
                title, size=12, bold=True, color=NAVY, line_spacing=1.0)


def _draw_project_page(slide, x, y, w, h):
    _panel_shell(slide, x, y, w, h,
                 "artifact · projects/tnf-paper.md",
                 "Project page — TNF paper")
    bx = x + Inches(0.25)
    bw = w - Inches(0.5)
    by = y + Inches(0.62)

    # frontmatter block (code-style)
    fm_h = Inches(0.78)
    add_rect(slide, bx, by, bw, fm_h,
             fill=tint(LINE, 0.4), line_color=None)
    fm_lines = [
        "status: active",
        "goal_model: okr",
        "deadline: 2026-10-31 (Cell Reports)",
        "collaborators: Prof Smith, Ravi (postdoc)",
    ]
    for i, line in enumerate(fm_lines):
        add_textbox(slide, bx + Inches(0.1), by + Inches(0.06 + i * 0.17),
                    bw - Inches(0.2), Inches(0.17),
                    line, size=8, color=INK, line_spacing=1.0, font="Consolas")
    by += fm_h + Inches(0.12)

    # Checkpoints heading
    add_textbox(slide, bx, by, bw, Inches(0.22),
                "## Checkpoints",
                size=10, bold=True, color=NAVY, line_spacing=1.0,
                font="Consolas")
    by += Inches(0.26)

    checkpoints = [
        (GREEN, "Methods drafted",  "Apr 15", "DoD: 2000 words, cites done"),
        (AMBER, "Results drafted",  "May 15", "DoD: numbers + structure"),
        (MUTED, "Figures v3",       "May 31", "DoD: 6 of 6 finalised"),
        (MUTED, "Co-author review", "Jun 15", "DoD: all replies incorporated"),
    ]
    for color, name, date, dod in checkpoints:
        add_dot(slide, bx + Inches(0.08), by + Inches(0.09),
                Inches(0.1), color)
        add_textbox(slide, bx + Inches(0.22), by,
                    Inches(1.6), Inches(0.2),
                    name, size=9, bold=True, color=NAVY, line_spacing=1.0)
        add_textbox(slide, bx + Inches(1.9), by,
                    Inches(0.75), Inches(0.2),
                    date, size=8, color=MUTED, line_spacing=1.0)
        add_textbox(slide, bx + Inches(2.7), by,
                    bw - Inches(2.7), Inches(0.2),
                    dod, size=8, color=MUTED, line_spacing=1.0)
        by += Inches(0.22)


def _draw_oneonone(slide, x, y, w, h):
    _panel_shell(slide, x, y, w, h,
                 "artifact · /one-on-one sara",
                 "1:1 agenda draft — Sara (postdoc) · Apr 24")
    bx = x + Inches(0.25)
    bw = w - Inches(0.5)
    by = y + Inches(0.62)

    sections = [
        ("She owes me  (2)",
         [("Organoid protocol v2", "committed Apr 8 · due Apr 22"),
          ("Figure 4 first draft",  "committed Apr 14 · due Apr 30")]),
        ("I owe her  (1)",
         [("Feedback on cell type B results", "due Apr 23")]),
        ("Recent progress",
         [("Media v2 reliable across 6 passages", ""),
          ("Cell type B contamination — root cause traced to incubator", "")]),
        ("Open blockers",
         [("Antibody reorder pending — ETA Apr 28", "")]),
    ]
    for title_txt, items in sections:
        add_textbox(slide, bx, by, bw, Inches(0.2),
                    title_txt, size=9, bold=True, color=TEAL, line_spacing=1.0)
        by += Inches(0.2)
        for item, meta in items:
            add_textbox(slide, bx + Inches(0.15), by,
                        Inches(3.2), Inches(0.18),
                        f"\u2022  {item}",
                        size=8, color=INK, line_spacing=1.15)
            if meta:
                add_textbox(slide, bx + Inches(3.4), by,
                            bw - Inches(3.4), Inches(0.18),
                            meta,
                            size=7, color=MUTED, line_spacing=1.15, font="Consolas")
            by += Inches(0.19)
        by += Inches(0.05)


def _draw_collab_report(slide, x, y, w, h):
    _panel_shell(slide, x, y, w, h,
                 "artifact · /report tnf for prof-smith",
                 "Weekly report — TNF paper · for Prof Smith")
    bx = x + Inches(0.25)
    bw = w - Inches(0.5)
    by = y + Inches(0.62)

    add_textbox(slide, bx, by, bw, Inches(0.2),
                "Week of Apr 14 — Apr 20",
                size=9, color=MUTED, line_spacing=1.0)
    by += Inches(0.23)

    # This week
    add_textbox(slide, bx, by, bw, Inches(0.22),
                "Progress this week",
                size=10, bold=True, color=NAVY, line_spacing=1.0)
    by += Inches(0.22)
    for line in [
        "Methods section drafted (2000 words)",
        "Figures 1–3 finalised at v3",
        "Replicate run matches primary data within 5%",
    ]:
        add_textbox(slide, bx + Inches(0.15), by,
                    bw - Inches(0.15), Inches(0.19),
                    f"\u2713  {line}",
                    size=8, color=INK, line_spacing=1.1)
        by += Inches(0.2)
    by += Inches(0.08)

    # Next week
    add_textbox(slide, bx, by, bw, Inches(0.22),
                "Next week",
                size=10, bold=True, color=NAVY, line_spacing=1.0)
    by += Inches(0.22)
    for line in [
        "Results section drafting (target: Fri)",
        "Figures 4–6 first pass",
    ]:
        add_textbox(slide, bx + Inches(0.15), by,
                    bw - Inches(0.15), Inches(0.19),
                    f"\u2192  {line}",
                    size=8, color=INK, line_spacing=1.1)
        by += Inches(0.2)
    by += Inches(0.06)

    # Asks
    add_textbox(slide, bx, by, bw, Inches(0.22),
                "Your input needed",
                size=10, bold=True, color=AMBER, line_spacing=1.0)
    by += Inches(0.22)
    add_textbox(slide, bx + Inches(0.15), by,
                bw - Inches(0.15), Inches(0.4),
                "\u2022  Methods section comments before Results starts?\n"
                "\u2022  Cell Reports 6-panel figure limit — stick to it?",
                size=8, color=INK, line_spacing=1.3)
    by += Inches(0.44)

    # output variants footer
    add_line(slide, bx, by, bx + bw, by, color=LINE, width_pt=0.5)
    by += Inches(0.08)
    add_textbox(slide, bx, by, bw, Inches(0.2),
                "Saved to exports/  ·  .md  ·  .email.md  ·  .slides.md",
                size=8, bold=True, color=TEAL, line_spacing=1.0, font="Consolas")


def _draw_duration_prompt(slide, x, y, w, h):
    _panel_shell(slide, x, y, w, h,
                 "prompt · duration learning",
                 "How long does this take? (learned from history)")
    bx = x + Inches(0.25)
    bw = w - Inches(0.5)
    by = y + Inches(0.7)

    # tool bubble
    bubble_h = Inches(1.15)
    add_rrect(slide, bx, by, bw, bubble_h,
              fill=tint(TEAL, 0.92),
              line_color=TEAL, line_width=Pt(0.5))
    add_dot(slide, bx + Inches(0.2), by + Inches(0.2),
            Inches(0.16), TEAL)
    add_textbox(slide, bx + Inches(0.4), by + Inches(0.1),
                Inches(1.5), Inches(0.22),
                "Claude",
                size=9, bold=True, color=TEAL, line_spacing=1.0)
    add_textbox(slide, bx + Inches(0.25), by + Inches(0.32),
                bw - Inches(0.5), bubble_h - Inches(0.4),
                "Based on your last 8 figure drafts (median 2h 20m), "
                "I\u2019ll block 2.5h tomorrow 09:00–11:30 in your "
                "protected writing block.\n\nConfirm, or give a different length.",
                size=9.5, color=INK, line_spacing=1.3)
    by += bubble_h + Inches(0.15)

    # user reply
    reply_w = Inches(2.4)
    rx = bx + bw - reply_w
    add_rrect(slide, rx, by, reply_w, Inches(0.38),
              fill=NAVY, line_color=None)
    add_textbox(slide, rx + Inches(0.2), by + Inches(0.08),
                reply_w - Inches(0.4), Inches(0.22),
                "yes — confirm",
                size=10, color=WHITE, line_spacing=1.0, font="Consolas")
    by += Inches(0.53)

    # data source
    add_textbox(slide, bx, by, bw, Inches(0.2),
                "Underlying data  ·  meta/task-history.json",
                size=8, bold=True, color=TEAL, line_spacing=1.0)
    by += Inches(0.2)
    add_textbox(slide, bx, by, bw, Inches(0.4),
                "figure_draft: { median_min: 140, p90_min: 210, n: 8 }",
                size=8, color=INK, line_spacing=1.2, font="Consolas")


def build_slide_6(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(slide, 6, "Four things the tool produces",
               eyebrow="Outputs")

    add_textbox(slide, MARGIN, Inches(1.42),
                Inches(12.33), Inches(0.5),
                "A project page, a 1:1 agenda draft, a weekly report for a collaborator, and a time-estimate prompt. All come from the same markdown files. Mockups below.",
                size=12, color=MUTED, line_spacing=1.2)

    # 2x2 grid
    gutter_x = Inches(0.3)
    gutter_y = Inches(0.3)
    grid_x0 = MARGIN
    grid_y0 = Inches(2.0)
    card_w = (SLIDE_W - MARGIN * 2 - gutter_x) / 2
    card_h = (Inches(4.95) - gutter_y) / 2

    _draw_project_page(slide, grid_x0, grid_y0, card_w, card_h)
    _draw_oneonone(slide, grid_x0 + card_w + gutter_x, grid_y0, card_w, card_h)
    _draw_collab_report(slide, grid_x0, grid_y0 + card_h + gutter_y, card_w, card_h)
    _draw_duration_prompt(slide, grid_x0 + card_w + gutter_x,
                          grid_y0 + card_h + gutter_y, card_w, card_h)


# ---------- main ----------

def build(out_name="PRESENTATION.pptx"):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_slide_1(prs)
    build_slide_2(prs)
    build_slide_3(prs)
    build_slide_4(prs)
    build_slide_5(prs)
    build_slide_6(prs)

    out = Path(__file__).resolve().parent.parent / out_name
    prs.save(out)
    print(f"Wrote {out}")
    return out


if __name__ == "__main__":
    name = sys.argv[1] if len(sys.argv) > 1 else "PRESENTATION.pptx"
    build(name)
